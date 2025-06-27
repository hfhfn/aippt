import datetime
import logging
import os
import json
import random
import glob

import pytz
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from llm import chat
from langchain.schema import HumanMessage, AIMessage
from langchain_community.chat_message_histories import ChatMessageHistory


# 配置日志记录器
logging.basicConfig(
    level=logging.DEBUG,  # 设置日志级别为DEBUG
    format='%(asctime)s - %(levelname)s - %(message)s',  # 设置日志格式
    handlers=[
        # logging.FileHandler("aippt.log"),  # 将日志写入文件
        logging.StreamHandler()  # 将日志输出到控制台
    ]
)


# 初始化对话历史记录
chat_history = ChatMessageHistory()


base_dir = os.path.abspath(os.path.dirname(__file__))
# 缓存目录
cache_dir = os.path.join(base_dir, "../output/temp/ppt")
os.makedirs(cache_dir, exist_ok=True)
# 输出目录
ppt_dir = os.path.join(base_dir, "../output/ppt")
os.makedirs(ppt_dir, exist_ok=True)


# 生成PPT内容
def generate_ppt_content(topic, pages):
    # 输出格式
    output_format = json.dumps({
        "title": "example title",
        "pages": [
            {
                "title": "title for page 1",
                "content": [
                    {
                        "title": "title for paragraph 1",
                        "description": "detail for paragraph 1",
                    },
                    {
                        "title": "title for paragraph 2",
                        "description": "detail for paragraph 2",
                    },
                ],
            },
            {
                "title": "title for page 2",
                "content": [
                    {
                        "title": "title for paragraph 1",
                        "description": "detail for paragraph 1",
                    },
                    {
                        "title": "title for paragraph 2",
                        "description": "detail for paragraph 2",
                    },
                    {
                        "title": "title for paragraph 3",
                        "description": "detail for paragraph 3",
                    },
                ],
            },
        ],
    }, ensure_ascii=True)

    # prompt
    prompt = f'''我要准备1个关于{topic}的PPT，要求一共写{pages}页，请你根据主题生成详细内容，不要省略。
                按这个JSON格式输出{output_format}，只能返回JSON，
                切记：1. JSON不要用```json```包裹，
                     2. 内容要用中文，
                     3. 每页字数不要超过250个字，
                     4. 标题前面不要写“第几页”'''

    # print(prompt)

    # 添加用户消息到对话历史
    chat_history.add_user_message(prompt)
    next_prompt = ""
    attempts = 0
    max_attempts = 5  # 最大尝试次数

    while attempts < max_attempts:
        attempts += 1

        # 构建消息列表
        if next_prompt:
            messages = [HumanMessage(content=next_prompt)]
        else:
            messages = [HumanMessage(content=prompt)]
        for message in chat_history.messages:
            if isinstance(message, HumanMessage):
                messages.append(HumanMessage(content=message.content))
            elif isinstance(message, AIMessage):
                messages.append(AIMessage(content=message.content))

        # 调用llm生成PPT内容
        output = chat.invoke(messages).content
        print(output)

        # 添加AI消息到对话历史
        chat_history.add_ai_message(output)

        try:
            if output.startswith("```json"):
                output = output.replace("```json", "").replace("```", "")
            ppt_content = json.loads(output)

            cache_files = glob.glob(os.path.join(cache_dir, "*.txt"))
            if len(cache_files) > 10:
                # 按修改时间排序，删除最旧的文件
                cache_files.sort(key=os.path.getmtime)
                os.remove(cache_files[0])

            with open(f"{cache_dir}/{topic}.txt", "w", encoding="utf-8") as f:
                json.dump(ppt_content, f, ensure_ascii=False, indent=4)
            return ppt_content
        except json.JSONDecodeError:
            print("生成的内容格式错误，重新生成...")
            next_prompt = f"生成的JSON格式错误，请重新按照如下提示：\n{prompt}\n生成符合格式的JSON内容。"
            continue
    print("尝试次数过多，生成失败！")
    return None

def generate_ppt_file(topic, ppt_content, design_number, layout_index):
    """生成PPT文件

    Args:
        topic: PPT主题
        ppt_content: PPT内容字典，包含title和pages
        design_number: 设计模板编号
        layout_index: 布局索引

    Returns:
        str: 生成的PPT文件路径
    """
    logging.info(f"开始生成PPT文件，主题：{topic}，设计模板：{design_number}，布局索引：{layout_index}")

    if ppt_content is None:
        print("PPT内容生成失败，请重新尝试！")
        return "PPT内容生成失败，请重新尝试！"

    # 1. 初始化PPT对象
    ppt = initialize_presentation(design_number)

    # 2. 添加首页
    add_title_slide(ppt, ppt_content['title'])

    # 3. 处理内容页
    process_content_slides(ppt, ppt_content['pages'], design_number, layout_index)

    # 4. 保存文件
    ppt_path = save_presentation(ppt, topic)
    return ppt_path


def initialize_presentation(design_number):
    """初始化PPT对象"""
    template_path = f"Designs/Design-{design_number}.pptx"
    if os.path.exists(template_path):
        return Presentation(template_path)
    template_path = f"Designs/Design-{design_number}.potx"
    if os.path.exists(template_path):
        return Presentation(template_path)
    logging.info(f"模板文件 {template_path} 不存在，将使用空白PPT模板")

    return Presentation()


def add_title_slide(ppt, title):
    """添加标题页"""
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])  # title&subtitle layout

    # 安全地设置标题
    title_placeholder = None
    subtitle_placeholder = None

    # 查找标题和副标题占位符
    for ph in slide.placeholders:
        try:
            ph_type = ph.placeholder_format.type
            if ph_type == 1:  # 标题占位符
                title_placeholder = ph
            elif ph_type == 2:  # 副标题占位符
                subtitle_placeholder = ph
        except Exception as e:
            logging.info(f"检查占位符时出错: {e}")
            continue

    # 如果没有找到标准的标题占位符，尝试使用索引方式（向后兼容）
    if title_placeholder is None:
        try:
            if len(slide.placeholders) > 0:
                title_placeholder = slide.placeholders[0]
        except (KeyError, IndexError) as e:
            logging.warning(f"无法找到标题占位符: {e}")

    if subtitle_placeholder is None:
        try:
            if len(slide.placeholders) > 1:
                subtitle_placeholder = slide.placeholders[1]
        except (KeyError, IndexError) as e:
            logging.info(f"无法找到副标题占位符: {e}")

    # 设置文本
    if title_placeholder:
        try:
            title_placeholder.text = title
        except Exception as e:
            logging.warning(f"设置标题失败: {e}")
    else:
        logging.warning("未找到可用的标题占位符")

    if subtitle_placeholder:
        try:
            subtitle_placeholder.text = "AI模型研究所"
        except Exception as e:
            logging.info(f"设置副标题失败: {e}")
    else:
        logging.info("未找到副标题占位符，跳过副标题设置")


def process_content_slides(ppt, pages, design_number, layout_index):
    """处理所有内容页"""
    logging.info(f'总共{len(pages)}页...')

    # 确定可用布局
    available_layouts = determine_available_layouts(ppt, layout_index)

    last_used_layout = -1

    for i, page in enumerate(pages):
        logging.info(f'生成第{i + 1}页:{page["title"]}')

        if design_number == 0:
            add_simple_content_slide(ppt, page)
            continue

        add_designed_content_slide(ppt, page, available_layouts, last_used_layout, i)


def determine_available_layouts(ppt, layout_index):
    """确定可用布局"""
    layout_count = len(ppt.slide_layouts)
    logging.info(f'当前ppt模板的布局数量为:{layout_count}')

    if layout_index in range(1, layout_count):
        return layout_index  # 可用布局索引, 一般ppt有 0-11 的布局，0给首页
    if layout_index == 0:
        return range(1, layout_count)
    return [1, 7, 8]  # 手动指定布局索引组合，layout_index=-1时使用


def add_simple_content_slide(ppt, page):
    """添加简单内容页(design_number=0时使用)"""
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # title&content layout

    # 设置标题
    slide.placeholders[0].text = page['title']

    # 添加正文内容
    content_placeholder = slide.placeholders[1]
    for sub_content in page['content']:
        print(sub_content)
        # 一级正文
        sub_title = content_placeholder.text_frame.add_paragraph()
        sub_title.text, sub_title.level = sub_content['title'], 1
        # sub_title.font.name = "微软雅黑"
        # sub_title.font.size = Pt(12)
        # 二级正文
        sub_description = content_placeholder.text_frame.add_paragraph()
        sub_description.text, sub_description.level = sub_content['description'], 2

    # 清理空占位符
    clean_empty_placeholders(slide)


def add_designed_content_slide(ppt, page, available_layouts, last_used_layout, slide_index):
    """添加设计内容页"""
    slide_added = False
    attempts = set()
    max_attempts = set(available_layouts) if isinstance(available_layouts, list) else set(range(1, len(ppt.slide_layouts)))

    while not slide_added and not attempts == max_attempts:

        # 选择布局
        layout_index = select_layout(available_layouts, last_used_layout, attempts)
        attempts.add(layout_index)

        try:
            # 获取布局
            if layout_index >= len(ppt.slide_layouts):
                logging.info(f"布局索引{layout_index}超出范围，最大索引为{len(ppt.slide_layouts) - 1}")
                continue

            slide_layout = ppt.slide_layouts[layout_index]

            # 检查占位符
            title_ph = find_title_placeholders(slide_layout)
            logging.info(f"标题占位符：{title_ph}")

            content_ph, content_parts = find_optimal_placeholder(slide_layout, page['content'])
            logging.info(f"内容占位符：{content_ph}")

            if not (title_ph and all(content_ph)):  # 保证标题和内容占位符都存在
                logging.info(f"布局{layout_index}缺少标题或内容占位符")
                if isinstance(available_layouts, int):
                    available_layouts = max_attempts
                continue

            # 添加幻灯片
            slide = ppt.slides.add_slide(slide_layout)
            slide_added = True
            last_used_layout = layout_index

            # 设置标题
            set_placeholder_text(slide, title_ph, page['title'])

            # 设置内容
            logging.info(f"第{slide_index + 1}页PPT, 使用了布局{layout_index}")
            fill_content_placeholder(slide, content_ph, content_parts)

            # 处理其他占位符
            process_additional_placeholders(slide, slide_index)

        except Exception as e:
            print(f"添加幻灯片时出错: {e}")
            continue


def select_layout(available_layouts, last_used_layout, attempts):
    """选择布局"""
    if isinstance(available_layouts, int):
        return available_layouts

    unused_layouts = list(set(available_layouts) - attempts)
    layout_index = random.choice(unused_layouts)
    # 确保新布局与上次不同
    while layout_index == last_used_layout and len(unused_layouts) > 1:
        layout_index = random.choice(unused_layouts)
    return layout_index


def find_title_placeholders(slide_layout):
    """查找标题占位符"""
    title_ph = None

    for ph in slide_layout.placeholders:
        if ph.placeholder_format.type == 1:  # 标题
            title_ph = ph

    return title_ph


def normalize_placeholder_type(ph_type):
    """将占位符类型统一转换为整数"""
    if isinstance(ph_type, PP_PLACEHOLDER_TYPE):
        return ph_type.value
    return int(ph_type)  # 兼容旧版整型


def find_optimal_placeholder(slide_layout, page_content):
    """基于文本内容智能选择最佳占位符"""
    text_content = page_dict2str(page_content)

    # 文本分析
    is_chinese = any('\u4e00' <= char <= '\u9fff' for char in text_content)
    content_length = len(text_content)

    # 计算精确需求
    required_area = calculate_ideal_area(content_length, is_chinese)
    # aspect_ratio = 1.77  # 默认16:9比例

    # 候选占位符评分
    best_ph = None
    best_score = float('inf')

    for ph in slide_layout.placeholders:
        score = score_placeholder(ph, required_area)
        if score is not None and score < best_score:
            best_score = score
            best_ph = ph

    # 动态调整选中的占位符
    if best_ph:
        # 初始化第二个最佳占位符
        second_best_ph = None
        second_best_score = float('inf')
        best_ph_area = best_ph.width * best_ph.height
        # 如果required_area大于ph_area，寻找第二个占位符
        if required_area > best_ph_area:
            for ph in slide_layout.placeholders:
                score = score_placeholder(ph, required_area)
                if score is not None and score < second_best_score and ph != best_ph:
                    second_best_score = score
                    second_best_ph = ph

        # 如果找到了第二个占位符，确保较大的占位符放在列表前面
        if second_best_ph:
            if best_ph.width * best_ph.height < second_best_ph.width * second_best_ph.height:
                best_ph, second_best_ph = second_best_ph, best_ph
            best_ph = [best_ph, second_best_ph]
        else:
            best_ph = [best_ph]

        if len(best_ph) == 1:
            adjust_placeholder(best_ph[0], text_content)
            return best_ph, [page_content]
        else:
            # 判断 content_part_num 是奇数还是偶数
            content_part_num = len(page_content)
            if content_part_num % 2 == 1:  # 奇数
                # 取 2/3 并向下取整
                first_part_content = page_content[:int(2 * content_part_num / 3)]
                second_part_content = page_content[int(2 * content_part_num / 3):]
            else:  # 偶数
                first_part_content = second_part_content = page_content[:int(content_part_num / 2)]

            first_text_content = page_dict2str(first_part_content)
            second_text_content = page_dict2str(second_part_content)

            adjust_placeholder(best_ph[0], first_text_content)
            adjust_placeholder(best_ph[1], second_text_content)
            return best_ph, [first_part_content, second_part_content]

    return [], []


def score_placeholder(ph, required_area):
    ph_type = ph.placeholder_format.type
    # ph_idx = ph.placeholder_format.idx
    # print(f"Placeholder Index: {ph_idx}, Type: {ph_type}, Type Class: {type(ph_type)}")
    # if isinstance(ph_type, bool):
    #     print(f"Unexpected boolean type for placeholder index {ph_idx}: {ph_type}")
    #     continue
    ph_type = normalize_placeholder_type(ph_type)
    if ph_type not in [2, 7]:  # 仅考虑内容和富文本
        return None
    # 计算当前占位符参数
    ph_area = ph.width * ph.height
    # ph_ratio = ph.width / ph.height

    # 评分标准(权重可调)
    size_score = abs(ph_area - required_area) / required_area
    # ratio_score = abs(ph_ratio - aspect_ratio)
    # total_score = 0.7 * size_score + 0.3 * ratio_score
    total_score = size_score
    return total_score


def page_dict2str(page_content):
    """将页面数据转换为字符串"""
    text_content = ""
    for item in page_content:
        text_content += f"{item['title']}\n{item['description']}\n"
    return text_content


def adjust_placeholder(placeholder, text_content):
    """精确调整占位符属性和样式"""
    # 计算文本特征
    line_count = text_content.count('\n') + 1
    max_line_length = max(len(line) for line in text_content.split('\n'))
    is_chinese = any('\u4e00' <= char <= '\u9fff' for char in text_content)

    # 智能字体设置
    font_size = Pt(12)  # 基础大小
    if len(text_content) > 300:
        font_size = Pt(10)
    elif len(text_content) < 100:
        font_size = Pt(14)

    # 应用样式
    tf = placeholder.text_frame
    tf.clear()  # 清空现有内容
    p = tf.paragraphs[0]
    p.text = text_content
    p.font.size = font_size
    p.font.name = 'Microsoft YaHei' if is_chinese else 'Calibri'

    try:
        """
        MSO_AUTO_SIZE.NONE                # 0 - 不自动调整
        MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT   # 1 - 调整形状适应文字
        MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE   # 2 - 调整文字适应形状（你当前的值）
        """
        # 自动调整策略 - 使用枚举值
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT if line_count > 5 else MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.word_wrap = True if max_line_length > 30 else False
        # 边距设置
        tf.margin_left = int(0.3 * 360000)  # 0.3cm
        tf.margin_right = int(0.3 * 360000)
    except Exception as e:
        logging.info(f"执行失败 {str(e)}")


def calculate_ideal_area(content_length, is_chinese=True):
    """
    基于印刷标准的精确面积计算
    :param content_length: 字符总数
    :param is_chinese: 是否为中文文本
    :return: 所需面积(EMU²)
    """
    # 标准参数
    empty_factor = 2  # 文本框排版预留空白区域系数
    LINE_SPACING = 1.2  # 行距系数
    MARGIN = 0.5 * 360000  # 0.5cm边距(EMU)

    # 根据语言选择参数
    if is_chinese:
        chars_per_cm2 = 3.0  # 中文字符/cm²(12pt)
        char_width = 0.35 * 360000  # 单个中文字符宽度(0.35cm)
    else:
        chars_per_cm2 = 6.0  # 英文字符/cm²(12pt)
        char_width = 0.15 * 360000  # 平均英文字符宽度(0.15cm)

    # 计算基础面积
    base_area = (content_length / chars_per_cm2) * (360000 ** 2)

    # 计算理想宽度(基于字符流式布局)
    ideal_width = min(
        int((content_length ** 0.5) * char_width * 1.5 + 2 * MARGIN),  # 平方根近似
        int(25 * 360000)  # 最大宽度25cm
    )

    # 计算动态高度
    ideal_height = max(
        int(base_area / ideal_width * LINE_SPACING + 2 * MARGIN),
        int(2 * 360000)  # 最小高度2cm
    )

    return ideal_width * ideal_height * empty_factor


def set_placeholder_text(slide, placeholder_info, text):
    """设置占位符文本"""
    try:
        ph = slide.placeholders[placeholder_info.placeholder_format.idx]
        ph.text = text
    except Exception as e:
        logging.info(f"设置占位符文本失败: {e}")


def fill_content_placeholder(slide, placeholder_info, content_parts):
    """填充内容占位符"""
    try:
        if len(placeholder_info) == 1:
            _fill_content_placeholder(slide, placeholder_info[0], content_parts[0])
        else:
            for placeholder, content in zip(placeholder_info, content_parts):
                _fill_content_placeholder(slide, placeholder, content)

    except Exception as e:
        print(f"填充内容占位符失败: {e}")


def _fill_content_placeholder(slide, placeholder, content):
    ph = slide.placeholders[placeholder.placeholder_format.idx]
    for sub_content in content:
        print(sub_content)
        # 一级正文
        sub_title = ph.text_frame.add_paragraph()
        sub_title.text, sub_title.level = sub_content['title'], 1
        # 二级正文
        sub_description = ph.text_frame.add_paragraph()
        sub_description.text, sub_description.level = sub_content['description'], 2


def process_additional_placeholders(slide, slide_index):
    """处理其他类型的占位符"""
    for ph in slide.placeholders:
        try:
            ph_type = ph.placeholder_format.type

            if ph_type == 1:  # 标题
                continue
            elif ph_type == 2 or ph_type == 7:  # 内容
                pass
                # if not ph.text.strip():
                #     ph.element.getparent().remove(ph.element)
            elif ph_type == 3:  # 图片
                pass  # 可以在这里添加默认图片
            elif ph_type == 4:  # 表格
                pass  # 可以在这里添加默认表格
            elif ph_type == 5:  # 图表
                pass  # 可以在这里添加默认图表
            elif ph_type == 14:  # 日期
                shanghai_tz = pytz.timezone('Asia/Shanghai')
                ph.text = datetime.datetime.now(shanghai_tz).strftime("%Y-%m-%d %H:%M:%S")
            elif ph_type == 15:  # 页脚
                ph.text = "AI模型研究所"
            elif ph_type == 13:  # 幻灯片编号
                ph.text = str(slide_index + 1)

        except Exception as e:
            print(f"处理占位符{ph.placeholder_format.type}时出错: {e}")


def clean_empty_placeholders(slide):
    """清理空占位符"""
    placeholders_to_remove = []

    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            # 检查占位符内容是否为空
            if shape.text.strip() == "":  # 如果文本为空或仅包含空白字符
                print(f"检测到空占位符{ph_type}，准备删除")
                placeholders_to_remove.append(shape)

    # 逆序删除
    for ph in reversed(placeholders_to_remove):
        try:
            sp = ph._element
            sp.getparent().remove(sp)
        except Exception as e:
            print(f"删除空占位符时出错: {e}")


def save_presentation(ppt, topic):
    """保存PPT文件"""
    ppt_path = f'{ppt_dir}/{topic}.pptx'
    ppt.save(ppt_path)
    return ppt_path


if __name__ == '__main__':
    last_topic = ""
    last_pages = 0
    while True:
        # 输入需求
        topic = input('输入主题:')
        pages = int(input('输入页数:'))
        template_num = int(input('输入模板编号(0-8):'))
        layout_index = int(input('输入布局编号(-1-11):'))
        # 生成PPT内容
        if os.path.exists(f"{cache_dir}/{topic}.txt") and last_topic == topic and last_pages == pages:
            ppt_content = json.load(open(f"{cache_dir}/{topic}.txt", "r", encoding="utf-8"))
            print("从缓存中读取PPT内容...")
        else:
            ppt_content = generate_ppt_content(topic, pages)
        # 生成PPT文件
        generate_ppt_file(topic, ppt_content, template_num, layout_index)
        last_topic = topic
        last_pages = pages
